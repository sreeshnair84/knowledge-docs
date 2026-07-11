#!/usr/bin/env python3
"""
Extract plain text from every .md/.pdf/.docx/.pptx file under docs/ into
_meta/corpus.json, batched by index range so a single call doesn't time out
on a large repo.

Usage:
    python3 extract_corpus.py <start_index> <end_index> [--docs-dir docs]

Run repeatedly with increasing ranges (0 60, 60 120, ...) until the reported
"Corpus size so far" equals the total file count. Safe to re-run: already-
extracted files are skipped.

Requires: pypdf, python-docx, python-pptx (pip install --break-system-packages
pypdf python-docx python-pptx)
"""
import os
import sys
import json
import time
import argparse

try:
    from pypdf import PdfReader
except ImportError:
    PdfReader = None
try:
    from docx import Document
except ImportError:
    Document = None
try:
    from pptx import Presentation
except ImportError:
    Presentation = None

META_DIR = "_meta"
FILELIST_PATH = os.path.join(META_DIR, "_extract_filelist.json")
CORPUS_PATH = os.path.join(META_DIR, "corpus.json")


def extract_pdf(path, max_pages=8, timeout=8):
    if PdfReader is None:
        return "[ERROR: pypdf not installed]"
    text = []
    start = time.time()
    try:
        reader = PdfReader(path)
        for page in reader.pages[:max_pages]:
            if time.time() - start > timeout:
                text.append("[TRUNCATED-TIMEOUT]")
                break
            try:
                t = page.extract_text()
                if t:
                    text.append(t)
            except Exception:
                continue
    except Exception as e:
        return f"[ERROR extracting pdf: {e}]"
    return "\n".join(text)


def extract_docx(path):
    if Document is None:
        return "[ERROR: python-docx not installed]"
    try:
        d = Document(path)
        return "\n".join(p.text for p in d.paragraphs[:500])
    except Exception as e:
        return f"[ERROR extracting docx: {e}]"


def extract_pptx(path):
    if Presentation is None:
        return "[ERROR: python-pptx not installed]"
    try:
        p = Presentation(path)
        texts = []
        for slide in p.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    texts.append(shape.text_frame.text)
        return "\n".join(texts)
    except Exception as e:
        return f"[ERROR extracting pptx: {e}]"


def extract_md(path):
    try:
        with open(path, encoding="utf-8", errors="ignore") as f:
            return f.read()
    except Exception as e:
        return f"[ERROR: {e}]"


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("start", type=int)
    ap.add_argument("end", type=int)
    ap.add_argument("--docs-dir", default="docs")
    args = ap.parse_args()

    os.makedirs(META_DIR, exist_ok=True)

    if not os.path.exists(FILELIST_PATH):
        files = []
        for root, _dirs, fns in os.walk(args.docs_dir):
            for fn in fns:
                if fn.lower().endswith((".md", ".pdf", ".docx", ".pptx")):
                    files.append(os.path.join(root, fn))
        files.sort()
        with open(FILELIST_PATH, "w", encoding="utf-8") as f:
            json.dump(files, f)
    else:
        with open(FILELIST_PATH, encoding="utf-8") as f:
            files = json.load(f)

    if os.path.exists(CORPUS_PATH):
        with open(CORPUS_PATH, encoding="utf-8") as f:
            out = json.load(f)
    else:
        out = {}

    t0 = time.time()
    for path in files[args.start:args.end]:
        if path in out:
            continue
        low = path.lower()
        if low.endswith(".md"):
            text = extract_md(path)
        elif low.endswith(".pdf"):
            text = extract_pdf(path)
        elif low.endswith(".docx"):
            text = extract_docx(path)
        elif low.endswith(".pptx"):
            text = extract_pptx(path)
        else:
            continue
        out[path] = text

    with open(CORPUS_PATH, "w") as f:
        json.dump(out, f)

    print(
        f"Processed up to index {args.end}. Total files: {len(files)}. "
        f"Corpus size so far: {len(out)}. Elapsed: {time.time() - t0:.1f}s"
    )
    if len(out) >= len(files):
        print("All files extracted. Proceed to find_duplicates.py.")


if __name__ == "__main__":
    main()
